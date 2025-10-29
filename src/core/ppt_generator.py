from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_pptx(json_code):
    path = "data/output.pptx"
    ppt = Presentation()
    slide_layout = ppt.slide_layouts[6]  # blank layout

    for slide_data in json_code:
        slide = ppt.slides.add_slide(slide_layout)

        # --- Title box ---
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True

        title_p = title_tf.add_paragraph()
        title_p.text = slide_data.get('title', 'Untitled Slide')
        title_p.font.size = Pt(28)     # slightly smaller title font
        title_p.font.bold = True

        # Estimate how many lines title might take
        title_text = title_p.text
        approx_lines = max(1, len(title_text) // 40 + 1)  # assume ~40 chars per line
        line_top = 0.6 + 0.5 + 0.3 * approx_lines  # dynamic Y pos in inches

        # --- Divider line ---
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(0.8),
            top=Inches(line_top),
            width=Inches(8.4),
            height=Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        line.line.fill.background()

        # --- Content box (starts below line) ---
        content_box = slide.shapes.add_textbox(Inches(1), Inches(line_top + 0.3), Inches(8), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        content = slide_data.get('content', '') or slide_data.get('bullets', '')

        if isinstance(content, str):
            p = tf.add_paragraph()
            p.text = content
            p.font.size = Pt(16)
            p.line_spacing = 1.3
            p.space_after = Pt(8)

        elif isinstance(content, list):
            for bullet_item in content:
                if isinstance(bullet_item, str):
                    bullet_text = bullet_item
                elif isinstance(bullet_item, dict):
                    bullet_text = bullet_item.get(
                        'point',
                        bullet_item.get('bullet', bullet_item.get('item', str(bullet_item)))
                    )
                else:
                    bullet_text = str(bullet_item)

                p = tf.add_paragraph()
                p.text = f"• {bullet_text}"
                p.font.size = Pt(16)
                p.space_before = Pt(4)
                p.line_spacing = 1.2

    ppt.save(path)
    return path
